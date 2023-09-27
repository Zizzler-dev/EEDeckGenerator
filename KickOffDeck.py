from re import template
from docx import Document
from pptx import Presentation
import streamlit as st
import os
import io
import datetime
import pandas as pd
from docx.enum.table import WD_TABLE_ALIGNMENT
from io import BytesIO
from pptx.util import Inches

def remove_slide(presentation, slide_index):
    xml_slides = presentation.slides._sldIdLst
    xml_slides.remove(xml_slides[slide_index])



def find_replace_variables(presentation, variables):
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for variable_key, variable_value in variables.items():
                            if variable_key in run.text:
                                run.text = run.text.replace(variable_key, variable_value)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for variable_key, variable_value in variables.items():
                                    if variable_key in run.text:
                                        run.text = run.text.replace(variable_key, variable_value)

contributionTable = pd.read_csv('jfed contribution chart.csv')
contributionTable.columns = contributionTable.columns.str.strip()
expected_columns = ['Employee Amount', 'Spouse Amount', 'Child Amount']
if not any(col in contributionTable.columns for col in expected_columns):
    st.warning("The uploaded contribution table is missing some expected columns.")
    st.stop()

for col in expected_columns:
    if col in contributionTable.columns:
        contributionTable[col] = contributionTable[col].apply(lambda x: f"${x:,.2f}" if not pd.isna(x) else None)



#st.write(contributionTable)
today = datetime.datetime.now()
this_year = today.year 
dec_31 = datetime.date(this_year+1, 12, 31)

# Create a Streamlit app
st.title("ICHRA Employee Deck Generator")

Client = st.text_input('Client Legal Name')

OERange = st.date_input(
    "Select OE Range",
    ( ),
    today,
    dec_31,
)

#st.write(f"{OERange[0].strftime('%m/%d/%Y')} - {OERange[1].strftime('%m/%d/%Y')}")
DepContribution = st.checkbox('Dependent Contribution')
LeftOverFunds = st.checkbox('Leftover Funds')
MedicareEligible = st.checkbox('Medicare Eligible')

Contribution = st.checkbox('Upload Contribution Table?')

if Contribution:
    contributionTable = st.file_uploader("Upload Contribution Table:")


    if contributionTable is None:
        st.warning("Please upload a contribution table to continue.")
        st.stop()
    
    contributionTable = pd.read_csv(contributionTable)
    contributionTable.columns = contributionTable.columns.str.strip()
    expected_columns = ['Employee Amount', 'Spouse Amount', 'Child Amount']
    if not any(col in contributionTable.columns for col in expected_columns):
        st.warning("The uploaded contribution table is missing some expected columns.")
        st.stop()

    for col in expected_columns:
        if col in contributionTable.columns:
            contributionTable[col] = contributionTable[col].apply(lambda x: f"${x:,.2f}" if not pd.isna(x) else None)


    st.write(contributionTable)

ppt_template = Presentation("ICHRA101.pptx") 
checklist_template = Presentation("zizzl health - Employee Pre Enrollment Checklist.pptx")
if st.button('SUBMIT'):
    # Load your PowerPoint template with styling and layout
     # Replace with your template filename

    # # Define variables for text replacement
    variables = {
        "XXXXX": Client,
        "DATERANGE": f"{OERange[0].strftime('%m/%d/%Y')} - {OERange[1].strftime('%m/%d/%Y')}",
        "MEDICAREX": 'Once you have selected your plan you must email zizzl health (support@zizzlhealth.com) to report your carrier and plan information. ' if MedicareEligible else '',
        "DEPENDENTSX": 'Employees can receive additional contributions based on the number of dependents' if DepContribution else ''
    }

    table_variables = {}
    contributions = contributionTable['Employee Amount']

    
    for i in range(65):
        if i == 0:
            key = '0-CONT'
            value_index = 0  # Use the first value
        elif 15 <= i <= 64:
            key = f"{i}CONT"
            value_index = i - 14  # Adjust the index for the value
        else:
            continue  # Skip keys that are not needed

        table_variables[key] = str(contributions[value_index])

    #st.write(table_variables)

    if(not LeftOverFunds):
        remove_slide(ppt_template, 14)

    # # Iterate through slides and replace text
    find_replace_variables(ppt_template, variables) 
    find_replace_variables(ppt_template, table_variables) 

    # # Create a stream to save the modified presentation
    ppt_stream = BytesIO()
    ppt_template.save(ppt_stream)

    # Create a download button for the customized presentation
    st.download_button(
        label="Download Kick Off Deck",
        data=ppt_stream.getvalue(),
        file_name=f"EE-Deck-{Client}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

checklist_stream = BytesIO()
checklist_template.save(checklist_stream)

# Create a download button for the customized presentation
st.download_button(
    label="Download Pre Enrollment Checklist",
    data=checklist_stream.getvalue(),
    file_name=f"Pre-Enrollment-Checklist.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
